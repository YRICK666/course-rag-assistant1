from __future__ import annotations

import argparse
import hashlib
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from charset_normalizer import from_path
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from tqdm import tqdm

from ppt_converter import PptConversionError, convert_ppt_to_pptx
from retriever import (
    MATERIALS_DIR,
    OUTPUTS_DIR,
    encode_texts,
    ensure_output_dirs,
    get_embedding_model_name,
    get_collection,
    reset_collection,
)


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt"}
SELECTABLE_EXTENSIONS = SUPPORTED_EXTENSIONS | {".ppt"}
EXTRACTED_DIR = OUTPUTS_DIR / "extracted_text"
CHUNKS_JSONL = EXTRACTED_DIR / "chunks.jsonl"
MANIFEST_PATH = OUTPUTS_DIR / "index_manifest.json"


@dataclass
class TextBlock:
    text: str
    source_path: str
    file_name: str
    file_type: str
    page_number: int | None = None
    slide_number: int | None = None
    paragraph_number: int | None = None
    paragraph_end: int | None = None


def clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def read_pdf(path: Path, materials_dir: Path = MATERIALS_DIR) -> list[TextBlock]:
    blocks: list[TextBlock] = []
    reader = PdfReader(str(path))
    relative_path = str(path.relative_to(materials_dir))

    for page_index, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        if text:
            blocks.append(
                TextBlock(
                    text=text,
                    source_path=relative_path,
                    file_name=path.name,
                    file_type=".pdf",
                    page_number=page_index,
                )
            )
    return blocks


def read_pptx(path: Path, materials_dir: Path = MATERIALS_DIR) -> list[TextBlock]:
    blocks: list[TextBlock] = []
    presentation = Presentation(str(path))
    relative_path = str(path.relative_to(materials_dir))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text or "")
                if text:
                    parts.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [clean_text(cell.text or "") for cell in row.cells]
                    row_text = " | ".join(cell for cell in cells if cell)
                    if row_text:
                        parts.append(row_text)

        text = clean_text("\n".join(parts))
        if text:
            blocks.append(
                TextBlock(
                    text=text,
                    source_path=relative_path,
                    file_name=path.name,
                    file_type=".pptx",
                    slide_number=slide_index,
                )
            )
    return blocks


def read_docx(path: Path, materials_dir: Path = MATERIALS_DIR) -> list[TextBlock]:
    blocks: list[TextBlock] = []
    document = Document(str(path))
    relative_path = str(path.relative_to(materials_dir))
    paragraph_number = 0

    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text or "")
        if not text:
            continue
        paragraph_number += 1
        blocks.append(
            TextBlock(
                text=text,
                source_path=relative_path,
                file_name=path.name,
                file_type=".docx",
                paragraph_number=paragraph_number,
                paragraph_end=paragraph_number,
            )
        )

    for table in document.tables:
        for row in table.rows:
            cells = [clean_text(cell.text or "") for cell in row.cells]
            text = clean_text(" | ".join(cell for cell in cells if cell))
            if not text:
                continue
            paragraph_number += 1
            blocks.append(
                TextBlock(
                    text=text,
                    source_path=relative_path,
                    file_name=path.name,
                    file_type=".docx",
                    paragraph_number=paragraph_number,
                    paragraph_end=paragraph_number,
                )
            )

    return blocks


def read_txt(path: Path, materials_dir: Path = MATERIALS_DIR) -> list[TextBlock]:
    relative_path = str(path.relative_to(materials_dir))
    content: str | None = None

    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            content = path.read_text(encoding=encoding)
            break
        except UnicodeDecodeError:
            continue

    if content is None:
        best = from_path(path).best()
        content = str(best) if best is not None else ""

    blocks: list[TextBlock] = []
    current: list[str] = []
    paragraph_number = 0

    def flush() -> None:
        nonlocal paragraph_number, current
        text = clean_text("\n".join(current))
        if text:
            paragraph_number += 1
            blocks.append(
                TextBlock(
                    text=text,
                    source_path=relative_path,
                    file_name=path.name,
                    file_type=".txt",
                    paragraph_number=paragraph_number,
                    paragraph_end=paragraph_number,
                )
            )
        current = []

    for line in content.splitlines():
        if line.strip():
            current.append(line.strip())
        elif current:
            flush()

    if current:
        flush()

    return blocks


def load_file(path: Path, materials_dir: Path = MATERIALS_DIR) -> list[TextBlock]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf(path, materials_dir=materials_dir)
    if suffix == ".pptx":
        return read_pptx(path, materials_dir=materials_dir)
    if suffix == ".docx":
        return read_docx(path, materials_dir=materials_dir)
    if suffix == ".txt":
        return read_txt(path, materials_dir=materials_dir)
    return []


def merge_adjacent_paragraph_blocks(
    blocks: list[TextBlock],
    max_paragraphs: int = 6,
    max_chars: int = 1600,
    overlap_paragraphs: int = 1,
) -> list[TextBlock]:
    if not blocks or blocks[0].file_type not in {".docx", ".txt"}:
        return blocks

    merged: list[TextBlock] = []
    index = 0
    while index < len(blocks):
        selected: list[TextBlock] = []
        total_chars = 0
        cursor = index

        while cursor < len(blocks) and len(selected) < max_paragraphs:
            block = blocks[cursor]
            next_size = total_chars + len(block.text)
            if selected and next_size > max_chars:
                break
            selected.append(block)
            total_chars = next_size
            cursor += 1

        if not selected:
            selected = [blocks[index]]
            cursor = index + 1

        first = selected[0]
        last = selected[-1]
        merged.append(
            TextBlock(
                text="\n\n".join(block.text for block in selected),
                source_path=first.source_path,
                file_name=first.file_name,
                file_type=first.file_type,
                paragraph_number=first.paragraph_number,
                paragraph_end=last.paragraph_end or last.paragraph_number,
            )
        )

        next_index = cursor - overlap_paragraphs
        index = max(index + 1, next_index)

    return merged


def iter_material_files(materials_dir: Path = MATERIALS_DIR, files: Iterable[Path] | None = None) -> Iterable[Path]:
    if files is not None:
        return sorted(
            path
            for path in files
            if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
        )
    if not materials_dir.exists():
        return []
    return sorted(
        path
        for path in materials_dir.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def warn_unconverted_ppt_files(materials_dir: Path = MATERIALS_DIR) -> int:
    if not materials_dir.exists():
        return 0

    count = 0
    for path in sorted(materials_dir.rglob("*.ppt")):
        if any(candidate.is_file() for candidate in path.parent.glob(f"{path.stem}*.pptx")):
            continue
        count += 1
        print(f"[WARN] 跳过未转换的 PPT：{path}。请先转换为 .pptx 后再建库。")
    return count


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path

    counter = 1
    while True:
        candidate = path.with_name(f"{path.stem}_{counter}{path.suffix}")
        if not candidate.exists():
            return candidate
        counter += 1


def archive_original_ppt(source: Path, outputs_dir: Path) -> Path:
    archive_dir = outputs_dir / "archived_original_ppt"
    archive_dir.mkdir(parents=True, exist_ok=True)
    target = unique_path(archive_dir / source.name)
    source.replace(target)
    return target


def find_converted_pptx(path: Path, materials_dir: Path) -> Path | None:
    candidates = sorted(path.parent.glob(f"{path.stem}*.pptx"))
    for candidate in candidates:
        if candidate.is_file() and candidate.resolve().is_relative_to(materials_dir.resolve()):
            return candidate
    return None


def resolve_selected_files(selected_files: Iterable[str | Path], materials_dir: Path) -> list[Path]:
    resolved: list[Path] = []
    materials_root = materials_dir.resolve()
    seen: set[Path] = set()
    for item in selected_files:
        path = Path(item)
        if not path.is_absolute():
            path = materials_dir / path
        path = path.resolve()
        if path != materials_root and materials_root not in path.parents:
            raise ValueError(f"所选文件必须位于 materials 目录内：{item}")
        if path in seen:
            continue
        seen.add(path)
        resolved.append(path)
    return resolved


def prepare_selected_files_for_index(
    selected_files: Iterable[str | Path] | None,
    *,
    materials_dir: Path,
    outputs_dir: Path,
) -> tuple[list[Path] | None, dict]:
    conversion = {
        "success_count": 0,
        "failure_count": 0,
        "successes": [],
        "failures": [],
    }
    if selected_files is None:
        return None, conversion

    prepared: list[Path] = []
    for source in resolve_selected_files(selected_files, materials_dir):
        if not source.exists() or not source.is_file():
            conversion["failures"].append(
                {
                    "file_name": str(source),
                    "message": "所选文件不存在，已跳过。",
                }
            )
            conversion["failure_count"] += 1
            continue

        suffix = source.suffix.lower()
        if suffix in SUPPORTED_EXTENSIONS:
            prepared.append(source)
            continue
        if suffix != ".ppt":
            continue

        relative_source = source.relative_to(materials_dir).as_posix()
        try:
            existing = find_converted_pptx(source, materials_dir)
            if existing:
                target = existing
            else:
                target = unique_path(source.with_suffix(".pptx"))
                convert_ppt_to_pptx(source, target)
            archived = archive_original_ppt(source, outputs_dir)
        except (PptConversionError, OSError, ValueError) as exc:
            conversion["failure_count"] += 1
            conversion["failures"].append(
                {
                    "file_name": relative_source,
                    "message": f"转换失败，原始 PPT 已保留。{exc}",
                }
            )
            continue

        prepared.append(target)
        conversion["success_count"] += 1
        conversion["successes"].append(
            {
                "file_name": relative_source,
                "converted_pptx": target.relative_to(materials_dir).as_posix(),
                "archived_original_ppt": archived.relative_to(outputs_dir / "archived_original_ppt").as_posix(),
                "message": "已转换为 PPTX，原始 PPT 已归档。",
            }
        )

    return sorted(prepared), conversion


def split_text(text: str, chunk_size: int = 900, overlap: int = 120) -> list[str]:
    text = clean_text(text)
    overlap = max(0, min(overlap, chunk_size // 2))
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    sentence_parts = re.split(r"(?<=[。！？!?；;\.\n])", text)
    chunks: list[str] = []
    current = ""

    for part in sentence_parts:
        part = part.strip()
        if not part:
            continue

        while len(part) > chunk_size:
            head = part[:chunk_size]
            chunks.append(head.strip())
            part = part[chunk_size - overlap :]

        if len(current) + len(part) + 1 <= chunk_size:
            current = f"{current}\n{part}".strip()
        else:
            if current:
                chunks.append(current.strip())
                tail = current[-overlap:] if overlap > 0 else ""
                current = f"{tail}\n{part}".strip()
            else:
                current = part

    if current:
        chunks.append(current.strip())

    return [chunk for chunk in chunks if chunk]


def make_chunk_id(source_path: str, location: str, chunk_index: int, text: str) -> str:
    digest = hashlib.sha1(f"{source_path}|{location}|{chunk_index}|{text}".encode("utf-8")).hexdigest()
    return digest


def block_location(block: TextBlock) -> str:
    if block.page_number is not None:
        return f"page-{block.page_number}"
    if block.slide_number is not None:
        return f"slide-{block.slide_number}"
    if block.paragraph_number is not None:
        if block.paragraph_end and block.paragraph_end != block.paragraph_number:
            return f"paragraph-{block.paragraph_number}-{block.paragraph_end}"
        return f"paragraph-{block.paragraph_number}"
    return "unknown"


def build_chunks(
    chunk_size: int = 900,
    overlap: int = 120,
    materials_dir: Path = MATERIALS_DIR,
    files: Iterable[Path] | None = None,
) -> tuple[list[str], list[str], list[dict]]:
    ids: list[str] = []
    documents: list[str] = []
    metadatas: list[dict] = []
    files = list(iter_material_files(materials_dir=materials_dir, files=files))

    for path in tqdm(files, desc="Reading materials"):
        try:
            blocks = load_file(path, materials_dir=materials_dir)
        except Exception as exc:
            print(f"[WARN] Failed to read {path}: {exc}")
            continue

        blocks = merge_adjacent_paragraph_blocks(blocks)

        for block in blocks:
            chunks = split_text(block.text, chunk_size=chunk_size, overlap=overlap)
            for chunk_index, chunk in enumerate(chunks, start=1):
                location = block_location(block)
                chunk_id = make_chunk_id(block.source_path, location, chunk_index, chunk)
                metadata = {
                    "source_path": block.source_path,
                    "file_name": block.file_name,
                    "file_type": block.file_type,
                    "page_number": block.page_number,
                    "slide_number": block.slide_number,
                    "paragraph_number": block.paragraph_number,
                    "paragraph_start": block.paragraph_number,
                    "paragraph_end": block.paragraph_end,
                    "chunk_index": chunk_index,
                    "location": location,
                }
                metadata = {key: value for key, value in metadata.items() if value is not None}
                ids.append(chunk_id)
                documents.append(chunk)
                metadatas.append(metadata)

    return ids, documents, metadatas


def write_chunks_jsonl(
    ids: list[str],
    documents: list[str],
    metadatas: list[dict],
    chunks_jsonl: Path = CHUNKS_JSONL,
) -> None:
    chunks_jsonl.parent.mkdir(parents=True, exist_ok=True)
    with chunks_jsonl.open("w", encoding="utf-8") as file:
        for chunk_id, document, metadata in zip(ids, documents, metadatas):
            file.write(
                json.dumps(
                    {"id": chunk_id, "text": document, "metadata": metadata},
                    ensure_ascii=False,
                )
                + "\n"
            )


def write_manifest(
    file_count: int,
    chunk_count: int,
    model_name: str,
    *,
    materials_dir: Path = MATERIALS_DIR,
    outputs_dir: Path = OUTPUTS_DIR,
    manifest_path: Path = MANIFEST_PATH,
    selected_files: list[Path] | None = None,
) -> None:
    manifest = {
        "materials_dir": str(materials_dir),
        "outputs_dir": str(outputs_dir),
        "embedding_model": model_name,
        "file_count": file_count,
        "chunk_count": chunk_count,
        "supported_extensions": sorted(SUPPORTED_EXTENSIONS),
        "scope": "selected" if selected_files is not None else "all",
    }
    if selected_files is not None:
        manifest["selected_files"] = [path.relative_to(materials_dir).as_posix() for path in selected_files]
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    with manifest_path.open("w", encoding="utf-8") as file:
        json.dump(manifest, file, ensure_ascii=False, indent=2)


def remove_existing_chunks_for_files(collection, source_paths: Iterable[str]) -> None:
    for source_path in source_paths:
        try:
            collection.delete(where={"source_path": source_path})
        except Exception as exc:
            print(f"[WARN] Failed to remove old chunks for {source_path}: {exc}")


def build_index(
    chunk_size: int = 900,
    overlap: int = 120,
    batch_size: int = 32,
    embedding_model: str | None = None,
    reset: bool = False,
    materials_dir: str | Path | None = None,
    outputs_dir: str | Path | None = None,
    selected_files: list[str | Path] | None = None,
) -> dict:
    resolved_materials_dir = Path(materials_dir) if materials_dir is not None else MATERIALS_DIR
    resolved_outputs_dir = Path(outputs_dir) if outputs_dir is not None else OUTPUTS_DIR
    extracted_dir = resolved_outputs_dir / "extracted_text"
    chunks_jsonl = extracted_dir / "chunks.jsonl"
    manifest_path = resolved_outputs_dir / "index_manifest.json"

    ensure_output_dirs(outputs_dir=resolved_outputs_dir)
    resolved_materials_dir.mkdir(parents=True, exist_ok=True)
    extracted_dir.mkdir(parents=True, exist_ok=True)

    model_name = embedding_model or get_embedding_model_name()
    prepared_selected_files, conversion_result = prepare_selected_files_for_index(
        selected_files,
        materials_dir=resolved_materials_dir,
        outputs_dir=resolved_outputs_dir,
    )
    skipped_ppt_count = 0 if prepared_selected_files is not None else warn_unconverted_ppt_files(resolved_materials_dir)
    files = list(iter_material_files(materials_dir=resolved_materials_dir, files=prepared_selected_files))
    ids, documents, metadatas = build_chunks(
        chunk_size=chunk_size,
        overlap=overlap,
        materials_dir=resolved_materials_dir,
        files=files,
    )

    if reset:
        collection = reset_collection(outputs_dir=resolved_outputs_dir)
    else:
        collection = get_collection(create=True, outputs_dir=resolved_outputs_dir)
        source_paths = [str(path.relative_to(resolved_materials_dir)) for path in files]
        remove_existing_chunks_for_files(collection, source_paths)

    if documents:
        for start in tqdm(range(0, len(documents), batch_size), desc="Embedding and indexing"):
            end = start + batch_size
            batch_docs = documents[start:end]
            batch_ids = ids[start:end]
            batch_metadatas = metadatas[start:end]
            embeddings = encode_texts(batch_docs, model_name=model_name, batch_size=batch_size)
            collection.upsert(
                ids=batch_ids,
                documents=batch_docs,
                metadatas=batch_metadatas,
                embeddings=embeddings,
            )

    write_chunks_jsonl(ids, documents, metadatas, chunks_jsonl=chunks_jsonl)
    write_manifest(
        file_count=len(files),
        chunk_count=len(documents),
        model_name=model_name,
        materials_dir=resolved_materials_dir,
        outputs_dir=resolved_outputs_dir,
        manifest_path=manifest_path,
        selected_files=files if selected_files is not None else None,
    )

    return {
        "mode": "reset" if reset else "update",
        "scope": "selected" if selected_files is not None else "all",
        "file_count": len(files),
        "requested_file_count": len(selected_files) if selected_files is not None else len(files),
        "chunk_count": len(documents),
        "embedding_model": model_name,
        "chroma_count": collection.count(),
        "skipped_ppt_count": skipped_ppt_count,
        "ppt_conversion": conversion_result,
        "indexed_files": [path.relative_to(resolved_materials_dir).as_posix() for path in files],
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Build local Chroma index for course materials.")
    parser.add_argument("--chunk-size", type=int, default=900)
    parser.add_argument("--overlap", type=int, default=120)
    parser.add_argument("--batch-size", type=int, default=32)
    parser.add_argument("--embedding-model", default=None)
    parser.add_argument(
        "--reset",
        action="store_true",
        help="Clear the old Chroma index in outputs before rebuilding from materials.",
    )
    args = parser.parse_args()

    summary = build_index(
        chunk_size=args.chunk_size,
        overlap=args.overlap,
        batch_size=args.batch_size,
        embedding_model=args.embedding_model,
        reset=args.reset,
    )
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
